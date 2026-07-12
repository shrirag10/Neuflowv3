"""Build docs/NeuFlow_v3_status.pptx — monochromatic academic status deck.

Regenerate after new results: edit constants and rerun.
"""

import sys, os, copy
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ---- monochromatic palette ----
INK    = RGBColor(0x1A, 0x1A, 0x1A)   # near-black: titles, body
ACCENT = RGBColor(0x33, 0x33, 0x33)   # dark gray: kickers, takeaways
MUTED  = RGBColor(0x8A, 0x8A, 0x8A)   # gray: footers, secondary
BOX_DK = RGBColor(0x33, 0x33, 0x33)   # flowchart: new/changed blocks
BOX_LT = RGBColor(0xEF, 0xEF, 0xEF)   # flowchart: unchanged blocks
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = 'Georgia'
OUT    = 'docs/NeuFlow_v3_status.pptx'


def add_text(s, x, y, w, h, text, size, bold=False, color=INK, line_spacing=1.0, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = line_spacing
        para.alignment = align
        r = para.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def header(s, kicker, title):
    add_text(s, 0.75, 0.42, 11.5, 0.3, kicker, 10.5, True, ACCENT)
    add_text(s, 0.73, 0.72, 11.9, 0.8, title, 24, True, INK)


def takeaway(s, text):
    return  # takeaway strips removed by request


def footer(s, num):
    add_text(s, 0.75, 7.02, 8.0, 0.30, 'NeuFlow v3, S. Srinivasan', 9, False, MUTED)
    add_text(s, 11.60, 7.02, 1.20, 0.30, str(num), 9, False, MUTED)


def box(s, x, y, w, h, text, dark=False, size=10, dashed=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.12
    shp.fill.solid()
    shp.fill.fore_color.rgb = BOX_DK if dark else BOX_LT
    shp.line.color.rgb = INK
    shp.line.width = Pt(1.0)
    if dashed:
        ln = shp.line._get_or_add_ln()
        d = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
        d.set('val', 'dash')
        shp.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(27432)
    tf.margin_top = tf.margin_bottom = Emu(9144)
    for i, line in enumerate(text.split('\n')):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = dark
        r.font.color.rgb = WHITE if dark else INK
    return shp


def arrow(s, x1, y1, x2, y2):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = INK
    conn.line.width = Pt(1.4)
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return conn



def framed_pic(s, path, x, y, w, h):
    s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    fr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fr.fill.background()
    fr.line.color.rgb = INK
    fr.line.width = Pt(1.2)
    fr.shadow.inherit = False


def scene_grid(s, panels):
    """panels: list of 4 (path, label) in reading order; wide 3.31:1 images."""
    W, H = 5.75, 1.74
    pos = [(0.75, 1.95), (6.85, 1.95), (0.75, 4.28), (6.85, 4.28)]
    for (path, label), (x, y) in zip(panels, pos):
        add_text(s, x, y - 0.32, W, 0.3, label, 11, True, INK)
        framed_pic(s, path, x, y, W, H)


def main():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    layout = p.slide_layouts[6]
    n = 0

    def slide():
        nonlocal n
        n += 1
        return p.slides.add_slide(layout), n

    # ================================================================ 1 · Title
    s, i = slide()
    add_text(s, 0.75, 2.5, 11.8, 1.0, 'NeuFlow v3', 40, True, INK)
    add_text(s, 0.75, 3.45, 11.8, 0.6, 'Queryable optical flow for edge devices', 18, False, ACCENT)
    add_text(s, 0.75, 4.35, 11.8, 1.0,
             'Status update. NeuFlow v2 backbone (frozen) + an implicit decoder that answers\n'
             'flow queries at arbitrary coordinates in O(N).\n\n'
             'Shriman Raghav Srinivasan, MS Robotics, Northeastern University, Field Robotics Lab', 13, False, MUTED)

    # ================================================== 2 · Problem, first principles
    s, i = slide()
    header(s, 'MOTIVATION', 'Why query optical flow?')
    add_text(s, 0.75, 1.7, 6.1, 4.4,
             'Optical flow tells you where each pixel in frame t moves\n'
             'to in frame t+1.\n\n'
             'Current networks always produce the full dense map at a\n'
             'fixed resolution. A lot of tasks never use most of it:\n\n'
             '•  registration wants a few hundred correspondences at\n'
             '    points it picks itself (corners, texture)\n'
             '•  sparse tracking only needs flow at feature points\n'
             '•  mosaicking needs matches in overlap regions, ideally\n'
             '    at sub-pixel positions\n\n'
             'On a small GPU, computing 479k values to use 800 of them\n'
             'is the difference between real time and not.', 13, False, INK)
    add_text(s, 7.2, 1.7, 5.4, 4.4,
             'Idea\n\n'
             'Keep everything in NeuFlow v2 that understands motion\n'
             '(backbone, matching, refinement), frozen.\n\n'
             'Swap only the last stage: instead of an upsampler that\n'
             'always produces the full map, use a decoder that returns\n'
             'flow at whatever coordinates you ask for.\n\n'
             'Cost then scales with the number of queries, O(N),\n'
             'instead of image area, O(HxW).', 13, False, INK)
    takeaway(s, 'The contribution is a new operating point — accuracy, compute, and resolution decoupled — not a leaderboard entry.')
    footer(s, i)

    # ================================================== 3 · Algorithm flow (flowchart)
    s, i = slide()
    header(s, 'METHOD', 'Pipeline: what stays from v2, what changes')

    bw, bh, gap = 1.48, 0.82, 0.28
    x0, yA, yB = 0.75, 2.05, 4.35
    chain = ['Image pair\n(t, t+1)', 'CNN backbone\nfeatures at 1/8, 1/16',
             'Cross-attention +\nglobal matching (1/16)', 'Recurrent refinement\n1 + 8 iterations',
             'Coarse flow\nat 1/8 scale']
    add_text(s, x0, yA - 0.42, 6.0, 0.3, 'NeuFlow v2', 11.5, True, INK)
    for k, t in enumerate(chain):
        box(s, x0 + k * (bw + gap), yA, bw, bh, t, dark=False, size=9.5)
        if k:
            arrow(s, x0 + k * (bw + gap) - gap + 0.02, yA + bh / 2, x0 + k * (bw + gap) - 0.02, yA + bh / 2)
    xe = x0 + 5 * (bw + gap)
    box(s, xe, yA, bw, bh, 'Convex upsampler\nfixed 8× grid', dark=True, size=9.5)
    arrow(s, xe - gap + 0.02, yA + bh / 2, xe - 0.02, yA + bh / 2)
    box(s, xe + bw + gap, yA, bw, bh, 'Dense flow map\nH × W, always', dark=False, size=9.5)
    arrow(s, xe + bw + 0.02, yA + bh / 2, xe + bw + gap - 0.02, yA + bh / 2)

    add_text(s, x0, yB - 0.42, 6.0, 0.3, 'NeuFlow v3 (this work)', 11.5, True, INK)
    for k, t in enumerate(chain):
        box(s, x0 + k * (bw + gap), yB, bw, bh, t, dark=False, size=9.5)
        if k:
            arrow(s, x0 + k * (bw + gap) - gap + 0.02, yB + bh / 2, x0 + k * (bw + gap) - 0.02, yB + bh / 2)
    box(s, xe, yB, bw, bh, 'Implicit decoder\n(next slide)', dark=True, size=9.5)
    arrow(s, xe - gap + 0.02, yB + bh / 2, xe - 0.02, yB + bh / 2)
    box(s, xe + bw + gap, yB, bw, bh, 'Flow at N requested\ncoordinates', dark=False, size=9.5)
    arrow(s, xe + bw + 0.02, yB + bh / 2, xe + bw + gap - 0.02, yB + bh / 2)
    box(s, xe - 0.1, yB + 1.15, bw + 0.2, 0.55, 'Queries (x, y) — any\ncontinuous positions', dark=False, size=9)
    arrow(s, xe + bw / 2, yB + 1.13, xe + bw / 2, yB + bh + 0.04)

    # brace: identical & frozen
    fr = box(s, x0 - 0.08, yB - 0.14, 5 * (bw + gap) - gap + 0.16, bh + 0.28, '', dashed=True)
    add_text(s, x0, yB + bh + 0.18, 6.4, 0.3, 'Identical to v2, weights frozen. matching quality inherited, not retrained',
             9.5, False, MUTED)
    add_text(s, xe - 0.05, yA + bh + 0.10, 3.6, 0.3, 'replaced stage (dark)', 9.5, False, MUTED)
    takeaway(s, 'One stage changes. Everything upstream of the 1/8-scale coarse flow is byte-identical to NeuFlow v2.')
    footer(s, i)

    # ================================================== 4 · Decoder flowchart
    s, i = slide()
    header(s, 'METHOD', 'Inside the decoder')
    bw2, bh2, gap2 = 2.05, 0.95, 0.42
    y = 2.6
    steps = [
        ('Query (x, y)\ncontinuous pixel\ncoordinate', False),
        ('Sample 3×3 windows\nfrom 4 feature maps:\ncontext · 1/8 · 1/16 ·\nwarped frame-1', False),
        ('Gated hierarchical\nfusion\n(InfiniDepth Eq. 3)', False),
        ('MLP → softmax over\n9 neighbors + bilinear\ncandidate (AnyFlow)', True),
        ('Convex blend of local\ncoarse-flow values\n= flow(x, y)', False),
    ]
    x = 0.75
    for k, (t, dark) in enumerate(steps):
        box(s, x, y, bw2, bh2, t, dark=dark, size=9.5)
        if k:
            arrow(s, x - gap2 + 0.04, y + bh2 / 2, x - 0.04, y + bh2 / 2)
        x += bw2 + gap2
    add_text(s, 0.75, 4.0, 11.9, 1.9,
             'Properties by construction, not by training\n\n'
             '•  Bounded output: the result is a weighted average of nearby coarse-flow values, so a bad weighting degrades\n'
             '    gracefully instead of producing garbage. (The earlier unbounded head never trained above its own init.)\n'
             '•  Exact init: bilinear interpolation is one particular weight setting, so the untrained decoder starts at 2.48 px EPE.\n'
             '•  The v2 upsampler is the fixed-grid special case of this mechanism, so v2 accuracy is reachable in principle.',
             12, False, INK)
    takeaway(s, 'The head predicts how to blend evidence, not the answer itself — the single design change that made training productive.')
    footer(s, i)

    # ============================================ 5 · Stage 0: untrained
    s, i = slide()
    header(s, 'RESULTS', 'Starting point: no training at all')
    scene_grid(s, [
        ('results/panels/scene_input.png', 'Input, VKITTI2 Scene18'),
        ('results/panels/scene_gt.png', 'Ground truth'),
        ('results/panels/scene_v2.png', 'NeuFlow v2, EPE 2.11 px'),
        ('results/panels/scene_v3_untrained.png', 'NeuFlow v3 untrained, EPE 2.26 px'),
    ])
    add_text(s, 0.75, 6.35, 11.9, 0.75, 'No decoder training at all: the bilinear-prior init already produces usable flow (2.48 px EPE over the full validation set,\n0.15 px behind v2), with querying functional and exact. Every training run afterwards had to beat this to be accepted.', 11.5, False, INK)
    footer(s, i)

    # ============================================ 6 · Stage 1: vkitti2
    s, i = slide()
    header(s, 'RESULTS', 'Trained on VKITTI2 only')
    scene_grid(s, [
        ('results/panels/scene_input.png', 'Input, VKITTI2 Scene18'),
        ('results/panels/scene_gt.png', 'Ground truth'),
        ('results/panels/scene_v2.png', 'NeuFlow v2, EPE 2.11 px'),
        ('results/panels/scene_v3_vkitti2.png', 'NeuFlow v3 (VKITTI2-trained), EPE 2.07 px'),
    ])
    add_text(s, 0.75, 6.35, 11.9, 0.75, 'Trained on 12,726 pairs: six VKITTI2 weather variants that share identical flow ground truth. Full-set result 2.39 px EPE,\nthe first checkpoint to improve on its own initialization, with no late-training collapse.', 11.5, False, INK)
    footer(s, i)

    # ============================================ 7 · Stage 2: chairs
    s, i = slide()
    header(s, 'RESULTS', 'Trained and evaluated on FlyingChairs')
    cw, chh = 2.92, 2.19
    xs = [0.75, 3.87, 6.99, 10.11]
    labels = ['Input, chairs val 00006', 'Ground truth', 'NeuFlow v2, EPE 1.42 px', 'NeuFlow v3 (chairs), EPE 1.35 px']
    paths = ['results/panels/chairs_input.png', 'results/panels/chairs_gt.png',
             'results/panels/chairs_v2.png', 'results/panels/chairs_v3.png']
    for x, lb, pt_ in zip(xs, labels, paths):
        add_text(s, x, 2.28, cw, 0.3, lb, 10.5, True, INK)
        framed_pic(s, pt_, x, 2.6, cw, chh)
    add_text(s, 0.75, 6.35, 11.9, 0.75, 'Across the 640 held-out chairs validation pairs: v2 2.24 px EPE (78.7% within 1 px), chairs-trained v3 2.40 px (76.6%).\nThe same model transfers to VKITTI2 driving scenes at 2.28 px, below v2 there, without ever seeing a road in training.', 11.5, False, INK)
    footer(s, i)

    # ============================================ 8 · Stage 3: mixed (best)
    s, i = slide()
    header(s, 'RESULTS', 'Trained on both datasets jointly (best model)')
    scene_grid(s, [
        ('results/panels/scene_input.png', 'Input, VKITTI2 Scene18'),
        ('results/panels/scene_gt.png', 'Ground truth'),
        ('results/panels/scene_v2.png', 'NeuFlow v2, EPE 2.11 px'),
        ('results/panels/scene_v3_mixed.png', 'NeuFlow v3 (mixed), EPE 2.06 px'),
    ])
    add_text(s, 0.75, 6.35, 11.9, 0.75, 'Both datasets sampled together in every batch (34,958 pairs), evaluated on held-out VKITTI2: 2.18 px EPE, 6% better than v2,\n76.4% within 1 px (v2: 77.6) and 3 px accuracy at parity. Sequential finetuning had failed at 2.50 px; joint sampling keeps both strengths.', 11.5, False, INK)
    footer(s, i)

    # ============================================ 8b · Head-to-head
    s, i = slide()
    header(s, 'RESULTS', 'v2 vs v3 on the same input')
    s.shapes.add_picture('results/visuals/head_to_head.png', Inches(1.35), Inches(1.55), width=Inches(10.6))
    footer(s, i)

    # ============================================ 8c · FlyingChairs examples
    s, i = slide()
    header(s, 'RESULTS', 'FlyingChairs examples')
    s.shapes.add_picture('results/visuals/chairs_examples.png', Inches(0.9), Inches(1.75), width=Inches(11.5))
    add_text(s, 0.9, 6.35, 11.5, 0.5,
             'Large, varied synthetic displacements. the motion statistics that taught the decoder its error-tail robustness.',
             11, False, MUTED)
    footer(s, i)

    # ============================================ 9 · Aggregate
    s, i = slide()
    header(s, 'RESULTS', 'Summary across training setups')
    s.shapes.add_picture('results/epe_by_regime.png', Inches(0.75), Inches(1.7), width=Inches(7.4))
    add_text(s, 8.5, 1.8, 4.2, 4.3,
             'Two informative negative results\n\n'
             'Sequential finetuning (chairs, then VKITTI2)\n'
             'reached 2.50 px. worse than either parent.\n'
             'Joint sampling was the remedy, confirmed by\n'
             'the mixed result.\n\n'
             'Fourier position encoding changed nothing\n'
             '(2.288 vs 2.275 px; 1 px accuracy identical).\n'
             'The sub-pixel gap is therefore not caused\n'
             'by missing positional information. one of\n'
             'two candidate explanations eliminated by a\n'
             'single controlled run.', 12, False, INK)
    takeaway(s, 'Each experiment either improved the model or eliminated a hypothesis; none was wasted compute.')
    footer(s, i)

    # ============================================ 10 · Efficiency, explained properly
    s, i = slide()
    header(s, 'EFFICIENCY', 'Sparse vs dense, and what each costs')
    add_text(s, 0.75, 1.62, 11.9, 1.45,
             'Dense output means the network computes flow for every pixel. 479,232 values for this input size. whether or not they are\n'
             'needed. This is the only mode NeuFlow v2 has: its upsampler is a convolution that always produces the full map in one 37 ms pass.\n\n'
             'Sparse output means computing flow only at N requested coordinates. Only v3 can do this: a 33 ms coarse pass shared by all queries,\n'
             'then 1.6 ms per batch of up to ~2,000 points. The sparse values are not approximations. they equal the dense output exactly.',
             12, False, INK)
    s.shapes.add_picture('results/latency_v2_v3.png', Inches(0.75), Inches(3.12), width=Inches(7.55))
    add_text(s, 8.55, 3.25, 4.1, 3.0,
             'Reading the charts\n\n'
             'Left: to get flow once, v2 dense\n'
             'and v3 sparse cost the same\n'
             '(~35 ms). v3 dense (327 ms) is an\n'
             'evaluation mode, not a use case.\n\n'
             'Right: once a pair is processed,\n'
             'each further question costs v3\n'
             '1.6 ms; v2 has no smaller unit of\n'
             'work than the full 37 ms frame.\n\n'
             'Video pipeline, 640x360 stream:\n'
             'v3 sparse-800:  63.6 FPS\n'
             'v2 dense:  60.3 FPS\n'
             'v3 + motion boxes:  47.1 FPS\n\n'
             'Parameters: 7.83 M (v3)\n'
             'vs 9.03 M (v2).', 11.5, False, INK)
    takeaway(s, 'v2 sells flow only by the full frame; v3 sells it by the question — at the same price for the first one and ~20× less for each after.')
    footer(s, i)

    # ============================================ 10b · Video throughput
    s, i = slide()
    header(s, 'EFFICIENCY', 'FPS on a real video stream')
    s.shapes.add_picture('results/fps_video.png', Inches(1.0), Inches(1.8), width=Inches(8.2))
    add_text(s, 9.5, 1.95, 3.2, 4.0,
             'Sixty frame pairs of a 640x360\n'
             'YouTube driving video, decoded,\n'
             'transferred, and processed -\n'
             'identical frames for every mode.\n\n'
             'v3 answering 800 targeted\n'
             'queries per pair outpaces v2\n'
             'producing its full map, and the\n'
             'live motion-detection mode\n'
             'holds 47 FPS. three times\n'
             'faster than a typical 15 FPS\n'
             'survey camera.', 12, False, INK)
    footer(s, i)

    # ============================================ 11 · Interface + GUI
    s, i = slide()
    header(s, 'INTERFACE', 'How to query the model')
    add_text(s, 0.75, 1.7, 6.2, 2.5,
             'A query is one continuous (x, y) coordinate; sub-pixel\n'
             'positions are valid inputs. N ranges from 1 to the full\n'
             'frame (479,232 at 384x1248). Decode cost is flat at\n'
             '1.6 ms up to roughly 2,000 queries per call.\n\n'
             'state = model.infer_coarse_state(img0, img1)   # once, 33 ms\n'
             'flow  = model.decode_queries(state, query_coords=q)\n'
             '        # q: [B, N, 2] pixel coords → [B, N, 2] flow', 12, False, INK)
    add_text(s, 0.75, 4.3, 6.2, 1.9,
             'Training configuration: batch size 4 (8 GB VRAM budget);\n'
             '4,096 supervision queries per image, half placed at motion\n'
             'boundaries; AdamW with a one-cycle schedule peaking at\n'
             '2e-4; backbone frozen throughout.', 12, False, INK)
    s.shapes.add_picture('results/visuals/query_gui_selftest.png', Inches(7.3), Inches(1.7), width=Inches(5.4))
    add_text(s, 7.3, 3.55, 5.4, 2.5,
             'Interactive tool (PyQt5, in the repository):\n'
             '•  Click any pixel for its flow; grid, boundary-adaptive,\n'
             '    and dense-overlay modes; CSV export.\n'
             '•  Region window: drag a rectangle. flow computed\n'
             '    only inside the selection.\n'
             '•  Video and YouTube sources with frame stepping and\n'
             '    real-time playback: live flow-based motion boxes at\n'
             '    47 FPS (640x360), ego-motion compensated.\n'
             '•  System-resources tab: live FPS, latency, GPU, VRAM,\n'
             '    CPU, RAM graphs. VRAM stays flat while interacting.', 10.5, False, INK)
    takeaway(s, 'The same API serves a robot asking for 800 correspondences and a human inspecting one pixel.')
    footer(s, i)

    # ============================================ 11b · GUI in action
    s, i = slide()
    header(s, 'INTERFACE', 'The GUI')
    s.shapes.add_picture('results/visuals/query_gui_selftest.png', Inches(0.75), Inches(1.85), width=Inches(6.0))
    add_text(s, 0.75, 3.85, 6.0, 0.9,
             'Adaptive queries over a dense overlay with the magnitude\n'
             'legend. several hundred answers from one cached state.', 10.5, False, MUTED)
    s.shapes.add_picture('results/visuals/query_gui_motion.png', Inches(7.0), Inches(1.85), width=Inches(6.0))
    add_text(s, 7.0, 3.85, 6.0, 0.9,
             'Playback with live motion detection: the translating region\n'
             'is boxed from the coarse flow at zero additional decode cost.', 10.5, False, MUTED)
    add_text(s, 0.75, 4.9, 11.8, 1.2,
             'A model selector switches between v3 and the v2 baseline in place: with v2, every interaction requires the full dense map\n'
             'to have been computed; with v3, the same interactions are answered from the cached coarse state in ~1.6 ms. The difference\n'
             'between the two architectures is directly felt in the tool.', 11.5, False, INK)
    footer(s, i)

    # ============================================ 11c · YouTube run proof
    s, i = slide()
    header(s, 'INTERFACE', 'Live run on public video: motion detection with measurements')
    s.shapes.add_picture('results/visuals/yt_run_viewer.png', Inches(0.75), Inches(1.7), width=Inches(6.3))
    s.shapes.add_picture('results/visuals/yt_run_resources.png', Inches(7.35), Inches(1.7), width=Inches(5.3))
    add_text(s, 0.75, 5.6, 11.8, 1.1,
             'Forty consecutive frame pairs of a public YouTube highway video (640x360), processed live in the GUI: mean 29.8 FPS\n'
             'end to end with motion boxes on the moving vehicles. The resource graphs were sampled during the same run; GPU memory\n'
             'stays flat at 769 MB because one cached backbone state serves every interaction.', 11, False, INK)
    footer(s, i)

    # ============================================ 12 · Objectives
    s, i = slide()
    header(s, 'OBJECTIVES', 'Where the three goals stand')
    add_text(s, 0.75, 1.85, 11.8, 4.2,
             'Goal 1, beat v2 accuracy.  Done: 2.18 vs 2.32 px mean EPE (6% better) with mixed training.\n'
             '1 px accuracy is within 1.2 points of v2 (76.4 vs 77.6) and 3 px is at parity.\n\n'
             'Goal 2, less compute, edge-viable.  Done for the sparse workload: same latency as a v2 full frame,\n'
             '13% fewer parameters, ~2.2 GB VRAM at inference, and repeat queries about 20x cheaper than v2\n'
             'recomputing. Dense-output mode stays slower; not the target use.\n\n'
             'Goal 3, do what v2 cannot.  Done: continuous-coordinate queries, output at any resolution, sparse\n'
             'matches dense exactly, plus a working interactive tool.',
             13, False, INK)
    takeaway(s, 'All three objectives are met or within measurement noise; remaining work is validation breadth, not capability.')
    footer(s, i)

    # ============================================ 13 · Next steps + asks
    s, i = slide()
    header(s, 'NEXT STEPS', 'Planned work and two requests')
    add_text(s, 0.75, 1.7, 6.1, 4.5,
             'Request 1. HPC access\n\n'
             'The Spring benchmark provides ground truth at twice the\n'
             'input resolution (1080p images, 4K flow). Evaluating there\n'
             'tests above-input-resolution querying. a question a\n'
             'fixed-resolution network cannot even accept. Its working\n'
             'set (hundreds of GB) and batch-scale training\n'
             'exceed the 8 GB laptop GPU this work has used so far.\n\n'
             'With HPC access, in order:\n'
             '•  Spring evaluation: v3 queried natively at 4K versus\n'
             '    v2 upsampled. the arbitrary-resolution claim, quantified\n'
             '•  FlyingThings3D pretraining (the standard second\n'
             '    second stage; 80 GB, batch 8–16)\n'
             '•  Longer mixed training at larger crops', 12, False, INK)
    add_text(s, 7.2, 1.7, 5.5, 4.5,
             'Request 2. Field Robotics Lab survey data\n\n'
             'The SeaBED AUV and UAV survey programs produce\n'
             'exactly the data this method serves:\n\n'
             '•  Sequential seafloor transect imagery (50–75% overlap,\n'
             '    slow motion, fine texture). registration and mosaicking\n'
             '    with sparse on-demand queries; the fine-motion regime\n'
             '    also stresses sub-pixel accuracy directly.\n'
             '•  Nadir UAV survey sequences with GPS/IMU tags -\n'
             '    georegistration provides quantitative pseudo ground truth.\n'
             '•  Vehicle navigation (DVL/INS) and camera calibration -\n'
             '    reprojection-based flow ground truth on static scenes,\n'
             '    no manual labeling.\n\n'
             'The ask is consecutive frames plus navigation data,\n'
             'not finished mosaics.', 12, False, INK)
    takeaway(s, 'Validation on public benchmarks is done on laptop hardware; both requests convert the method into thesis-scale evidence.')
    footer(s, i)

    # ============================================ 14 · FAQ
    s, i = slide()
    header(s, 'Q&A PREP', 'Questions I expect')
    add_text(s, 0.75, 1.65, 11.9, 4.5,
             'Why is flow defined at non-integer coordinates?  Bilinear interpolation makes feature maps continuous functions of position;\n'
             'an MLP composed with them is defined at every real (x, y). Integer pixels are the special case dense maps hard-code.\n\n'
             'Why freeze the backbone?  Joint training is a moving-target problem requiring ~800K steps at InfiniDepth scale; at a 30K budget\n'
             'the decoder diverges chasing shifting features (observed directly). Freezing also guarantees the v2 matching quality is inherited.\n\n'
             'Is sparse output an approximation?  No. It is the same function evaluated at fewer points; agreement with dense output is exact.\n\n'
             'Why did out-of-domain training beat in-domain training?  Mean EPE is dominated by the error tail. Diverse large motions train\n'
             'tail robustness; five driving scenes train memorization. Joint training then combined both, as the mixed result shows.\n\n'
             'What limits sub-pixel accuracy?  Not positional information. measured (PE ablation, null result). Remaining candidates: the\n'
             '1/8-scale coarse flow bounds recoverable detail, and large-motion training never supervises sub-pixel discrimination.\n\n'
             'Why batch size 4?  An 8 GB VRAM budget; the recipe is otherwise RAFT-standard and scales directly on larger GPUs.',
             11.5, False, INK, line_spacing=1.02)
    takeaway(s, 'Full derivations and the complete question list: docs/NeuFlow_v3_Report.md, sections 5–6.')
    footer(s, i)

    p.save(OUT)
    print(f'saved {OUT} with {n} slides')


if __name__ == '__main__':
    main()
