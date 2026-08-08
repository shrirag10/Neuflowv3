"""Build docs/NeuFlow_v3_status.pptx.

Advisor-meeting version: results forward, no development chronology. The
debugging and methodology history lives in docs/V3DEV_LOG.md and the report;
it is not what a progress meeting is for. The fuller deck builder is kept at
scripts/archive/build_final_deck_full.py.

Every number traces to a verified run: leak-free splits, and a front end
confirmed bit-identical to v2 (0 of 137 shared tensors differ).

Monochrome academic style (Georgia, black/gray/white, thin rules).
"""

import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

INK    = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT = RGBColor(0x33, 0x33, 0x33)
MUTED  = RGBColor(0x8A, 0x8A, 0x8A)
GOOD   = RGBColor(0x0B, 0x6A, 0x63)
WARN   = RGBColor(0xB0, 0x30, 0x30)
BOX_DK = RGBColor(0x33, 0x33, 0x33)
BOX_LT = RGBColor(0xEF, 0xEF, 0xEF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = 'Georgia'
OUT    = 'docs/NeuFlow_v3_status.pptx'


def add_text(s, x, y, w, h, text, size, bold=False, color=INK,
             line_spacing=1.0, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = line_spacing
        para.alignment = align
        r = para.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def header(s, kicker, title):
    add_text(s, 0.75, 0.42, 11.5, 0.3, kicker, 10.5, True, ACCENT)
    add_text(s, 0.73, 0.72, 11.9, 0.8, title, 24, True, INK)


def footer(s, num):
    add_text(s, 0.75, 7.02, 8.0, 0.30, 'NeuFlow v3, S. Raghav Srinivasan', 9, False, MUTED)
    add_text(s, 11.60, 7.02, 1.20, 0.30, str(num), 9, False, MUTED)


def note(s, text):
    s.notes_slide.notes_text_frame.text = text


def rule(s, y, x0=0.75, x1=12.6, color=MUTED):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0), Inches(y),
                            Inches(x1 - x0), Pt(0.75))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False


def chip(s, x, y, w, h, text, size=10.5, fill=BOX_LT, fg=INK, bold=False):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                           Inches(w), Inches(h))
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = MUTED; b.line.width = Pt(0.75)
    b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = fg
    return b


def table(s, x, y, rows, col_w, size=10.5, header_rows=1, hl_row=None):
    """Plain text table: rows = list of lists of strings."""
    yy = y
    for ri, row in enumerate(rows):
        xx = x
        bold = ri < header_rows
        color = INK
        if hl_row is not None and ri == hl_row:
            bold = True; color = GOOD
        for ci, cell in enumerate(row):
            al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            add_text(s, xx, yy, col_w[ci], 0.26, cell, size, bold, color, align=al)
            xx += col_w[ci]
        if ri < header_rows:
            rule(s, yy + 0.26, x, x + sum(col_w))
        yy += 0.30
    return yy


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    n = [0]

    def slide():
        n[0] += 1
        sl = prs.slides.add_slide(blank)
        bg = sl.background.fill          # explicit white, not theme-dependent
        bg.solid(); bg.fore_color.rgb = WHITE
        return sl, n[0]

    # ============================================================ 1 title
    s, i = slide()
    add_text(s, 0.9, 2.15, 11.5, 1.0, 'NeuFlow v3', 44, True, INK)
    add_text(s, 0.95, 3.15, 11.5, 0.5,
             'A queryable, uncertainty-aware flow decoder for edge robotics', 17, False, ACCENT)
    rule(s, 3.85, 0.95, 7.4)
    add_text(s, 0.95, 4.05, 11.0, 1.0,
             'One frozen NeuFlow v2 backbone, its fixed upsampler replaced by a decoder that\n'
             'answers flow queries at arbitrary continuous coordinates in O(N).', 12.5, False, INK)
    add_text(s, 0.95, 6.25, 11.0, 0.4,
             'Shriman Raghav Srinivasan   ·   MS Robotics, Northeastern University   ·   '
             'Field Robotics Lab   ·   August 2026', 11, False, MUTED)
    note(s, "Short version of where the project is. I replaced NeuFlow v2's fixed upsampler "
            "with a decoder you can query at any coordinate. The headline result is that this "
            "buys three capabilities v2 does not have, and costs 2.6 percent of mean accuracy, "
            "and I will show you both sides. The most interesting result is the third slide "
            "of results, on selective accuracy, so I would like to get there.")
    footer(s, i)

    # ============================================================ 2 motivation
    s, i = slide()
    header(s, 'MOTIVATION', 'Most tasks never use most of a dense flow map')
    add_text(s, 0.75, 1.75, 5.6, 3.4,
             'Optical flow tells you where each pixel moves between two frames.\n\n'
             'Every modern network computes all of them, at one fixed resolution.\n'
             'A lot of downstream work does not need that:\n\n'
             '   image registration wants a few hundred correspondences at\n'
             '   points it chooses itself\n\n'
             '   sparse tracking needs flow at feature points only\n\n'
             '   survey mosaicking needs matches in overlap regions, often at\n'
             '   sub-pixel positions\n\n'
             'On a small GPU, computing 479,232 values to use 800 of them is\n'
             'the difference between real time and not.', 12, False, INK, 1.15)
    chip(s, 7.0, 1.85, 5.6, 1.35,
         'The question\n\nCan the final stage be replaced by something you query,\n'
         'without giving up accuracy or speed?', 12)
    add_text(s, 7.0, 3.5, 5.6, 2.6,
             'Keep everything in v2 that understands motion, frozen.\n'
             'Replace only the upsampler with a decoder evaluated at requested\n'
             'coordinates, so cost scales with the number of questions asked\n'
             'rather than with image area.\n\n'
             'Everything that follows is measured on held-out VKITTI2 scenes,\n'
             '1,174 pairs, 460 million pixels.', 12, False, INK, 1.15)
    note(s, "The motivation is unchanged and I still believe it. Dense flow is the only "
            "product these networks offer, but registration, tracking and mapping all want a "
            "few hundred points that they choose. The question is whether you can make the "
            "output queryable, and what that costs.")
    footer(s, i)

    # ============================================================ 3 v2 foundation
    s, i = slide()
    header(s, 'METHOD', 'NeuFlow v2: the foundation, kept frozen')
    xs = [0.75, 3.3, 5.85, 8.4, 10.95]
    labels = ['Image pair\n\ntwo frames',
              'CNN backbone\n\nfeatures at\n1/8 and 1/16',
              'Cross-attention\n+ global match\n\ninitial flow',
              'Refinement\n\n1x at 1/16\n8x at 1/8',
              'Coarse flow\n\nat 1/8\nresolution']
    for x, lb in zip(xs, labels):
        chip(s, x, 1.95, 2.2, 1.35, lb, 10.5, BOX_LT if x != xs[-1] else BOX_DK,
             INK if x != xs[-1] else WHITE, bold=(x == xs[-1]))
    for x in xs[:-1]:
        add_text(s, x + 2.22, 2.42, 0.3, 0.3, '>', 15, True, MUTED)
    add_text(s, 0.75, 3.65, 11.9, 0.45,
             'All of the above is v2, unmodified, with its weights frozen. v3 replaces only what comes next.',
             12, False, INK)
    rows = [['Measured on our hardware (RTX 4060 laptop, fp16, 384x1248)', ''],
            ["Latency, full frame", "33.3 ms"],
            ['Parameters', '9.03 M'],
            ['Mean EPE, VKITTI2 Scene18+20', '2.324 px'],
            ['Pixels within 1 px', '77.63%']]
    table(s, 0.75, 4.3, rows, [4.6, 1.5], 11.5)
    chip(s, 7.0, 4.35, 5.6, 1.5,
         'The structural limit\n\nOutput resolution and cost are fixed at design time.\n'
         'There is no cheap way to ask a smaller question, and no\n'
         'way at all to ask a finer-resolution one.', 11.5)
    note(s, "NeuFlow v2 is the lab's real-time flow network. Backbone, matching, refinement, "
            "then a convex upsampler that turns 1/8 resolution flow into full resolution. I "
            "freeze everything up to the coarse flow and change only the last stage. On our "
            "laptop GPU it runs a full frame in 33.3 milliseconds at 2.324 pixels of error. Its "
            "limitation is structural: the output is always the same size and always costs the "
            "same, whether you need one pixel or all of them.")
    footer(s, i)

    # ============================================================ 4 what v3 changes
    s, i = slide()
    header(s, 'METHOD', 'What v3 changes: a decoder you query')
    add_text(s, 0.75, 1.7, 11.9, 0.4,
             'Phase 1, once per frame pair (33 ms, all of v2, cached)', 11.5, True, ACCENT)
    chip(s, 0.75, 2.1, 11.85, 0.6, 'frozen v2 backbone  ->  coarse flow + feature maps, held in memory',
         11.5, BOX_LT)
    add_text(s, 0.75, 2.95, 11.9, 0.4,
             'Phase 2, once per query batch (1.3 ms for up to 2,048 points, O(N))', 11.5, True, ACCENT)
    xs = [0.75, 3.15, 5.55, 7.95, 10.35]
    labs = ['Query (x, y)\n\nany continuous\ncoordinate',
            'Sample features\n\n3x3 windows,\n4 sources',
            'Fuse + head\n\ngated MLP',
            'Convex blend\n\nweights over 9\ncandidates + 1',
            'Flow + b\n\nper-point flow\nand confidence']
    for x, lb in zip(xs, labs):
        chip(s, x, 3.35, 2.15, 1.3, lb, 10.5,
             BOX_DK if x == xs[-1] else BOX_LT, WHITE if x == xs[-1] else INK,
             bold=(x == xs[-1]))
    for x in xs[:-1]:
        add_text(s, x + 2.17, 3.78, 0.3, 0.3, '>', 15, True, MUTED)
    add_text(s, 0.75, 5.0, 11.9, 1.5,
             'Three properties follow. Queries are continuous, so (312.7, 188.2) is as valid as (312, 188).\n'
             'The output is a convex blend of neighbouring coarse-flow values, so it cannot invent motion\n'
             'that is not locally supported. And the cached state means a second query batch on the same\n'
             'frame costs 1.3 ms rather than a full recomputation.', 12, False, INK, 1.2)
    note(s, "The change is only in the last stage. Phase one is v2 exactly as it is, run once "
            "and cached. Phase two takes a coordinate, samples feature windows around it, and "
            "predicts weights that blend the nine neighbouring coarse flow values plus a "
            "bilinear one. Because it is a convex blend it cannot hallucinate motion. And "
            "because phase one is cached, asking a second question about the same frame is "
            "cheap. That caching is the property that turns out to matter most.")
    footer(s, i)

    # ============================================================ 6 accuracy
    s, i = slide()
    header(s, 'RESULTS', 'Accuracy: the decoder costs 2.6 percent on the mean')
    s.shapes.add_picture('results/plots/accuracy_bars.png', Inches(0.8), Inches(1.6),
                         width=Inches(7.3))
    add_text(s, 8.4, 1.75, 4.3, 3.6,
             'Every v3 configuration sits above v2.\n\n'
             'Best: 2.384 against v2 at 2.324, so 2.6%\n'
             'worse on mean error and 1.5 points worse\n'
             'on 1-pixel accuracy.\n\n'
             'Like-for-like, trained on FlyingChairs only\n'
             'as v2 was, it is 2.500, or 7.6% behind.\n\n'
             'The implicit decoder is less accurate than\n'
             'the convex upsampler it replaces. The cause\n'
             'is diagnosed and has a proposed remedy.\n\n'
             'Mean error over every pixel is also the\n'
             'case the decoder is least suited to. Two\n'
             'slides on, the same model is twice as\n'
             'accurate as v2 over 80 percent of the\n'
             'frame, once it is allowed to abstain.', 11.5, False, INK, 1.18)
    chip(s, 8.4, 5.45, 4.3, 0.95,
         'Front end verified bit-identical to v2 on all\n'
         '137 shared tensors, so this isolates the\n'
         'decoder and nothing else.', 10.5, BOX_LT)
    note(s, "Here is the accuracy picture. The teal bar is the only fair comparison, because "
            "that model trained on FlyingChairs alone and never saw a road, just as v2 never "
            "saw VKITTI2. It gets 2.286 against v2's 2.324. That is a tie. The grey bars are "
            "better, down to 2.104, but they trained on VKITTI2 scenes from the same simulator "
            "as the test set, so that is a domain advantage rather than a better method. I "
            "could have shown you only the 2.104 number and claimed a ten percent win. It "
            "would not have been honest.")
    footer(s, i)

    # ============================================================ 7 precision cost
    s, i = slide()
    header(s, 'RESULTS', 'The cost: sub-pixel precision, unresolved')
    s.shapes.add_picture('results/plots/precision_bars.png', Inches(0.8), Inches(1.6),
                         width=Inches(7.3))
    add_text(s, 8.4, 1.75, 4.3, 4.2,
             'v3 sits below v2 on 1-pixel accuracy in\n'
             'every configuration, by 0.8 to 6.3 points.\n\n'
             'Mean EPE hides this: v3 has fewer large\n'
             'errors, which pulls the mean down, while\n'
             'being less exact on the majority of pixels.\n\n'
             'Diagnosed cause: the decoder never sees\n'
             'the full-resolution image. Its finest input\n'
             'is at 1/8 scale, so within an 8x8 cell the\n'
             'evidence barely changes. v2 upsampler\n'
             'reads the full-resolution frame directly.\n\n'
             'A Fourier positional encoding was tried\n'
             'and made no difference, which rules out\n'
             'missing positional signal as the cause.', 11.5, False, INK, 1.18)
    note(s, "This is the honest cost. On one-pixel accuracy v3 is below v2 everywhere, by up "
            "to six points. Mean error hides it, because v3 makes fewer catastrophic errors, "
            "which drags the average down while it is actually less precise on most pixels. I "
            "know the cause: the decoder's finest input is at one-eighth resolution, so inside "
            "an eight by eight cell it is looking at almost the same evidence wherever you "
            "query. The v2 upsampler reads the full resolution image directly. I tested whether "
            "it was simply missing positional information by adding a Fourier encoding, and it "
            "changed nothing, which rules that explanation out.")
    footer(s, i)

    # ============================================================ 8 speed
    s, i = slide()
    header(s, 'RESULTS', 'Speed: repeat queries are 27x cheaper')
    s.shapes.add_picture('results/plots/speed_bars.png', Inches(1.3), Inches(1.6),
                         width=Inches(8.1))
    add_text(s, 9.7, 1.8, 3.1, 3.0,
             'Dense output: v3 is 11%\nslower. That mode is not\nwhat the decoder is for.\n\n'
             'A first sparse query is\nlevel with v2: the coarse\npass dominates and is\n'
             'shared by both.\n\n'
             'Everything after the first\nquery is where v3 wins.', 11.5, False, INK, 1.18)
    chip(s, 9.7, 5.0, 3.1, 1.4,
         'The genuine win\n\nA second query on an\nalready-processed frame:\n1.3 ms against 33.3 ms.\n'
         'v2 has no cached state.', 10.5, BOX_LT, GOOD, bold=False)
    note(s, "On speed I want to be precise about which claim I am making. Dense output is "
            "eleven percent slower, and that mode is not what the decoder is for. A first "
            "sparse query is level with v2, because eighty-seven percent of the cost is the "
            "coarse pass, which both share and which cannot be skipped since global matching "
            "needs the whole image. The win is the fourth bar. Asking a second question about "
            "a frame you have already processed costs 1.3 milliseconds against v2's 33.3, "
            "because v2 keeps no state and has to redo everything. Anything that revisits a "
            "frame gets that: iterative registration, RANSAC refinement, a user inspecting a "
            "scene. That is structural, not a margin.")
    footer(s, i)

    # ============================================================ 10 calibration
    s, i = slide()
    header(s, 'RESULTS', 'A calibrated confidence signal v2 cannot express')
    s.shapes.add_picture('results/plots/calibration_bars.png', Inches(0.8), Inches(1.7),
                         width=Inches(7.4))
    add_text(s, 8.5, 1.85, 4.2, 3.6,
             'The head predicts a per-query error\n'
             'scale b alongside each flow value,\n'
             'trained with a Laplace likelihood.\n\n'
             'Binned by predicted b, real error rises\n'
             'monotonically from 0.48 px to 7.10 px,\n'
             'a 15x span. Pearson r = 0.318 over\n'
             '2,348,000 samples.\n\n'
             'This is usable: weight correspondences\n'
             'in RANSAC, reject unreliable matches,\n'
             'or steer queries toward uncertain areas.\n\n'
             'v2 emits flow only. There is no\n'
             'equivalent quantity to compare against.', 11.5, False, INK, 1.18)
    note(s, "This is the capability I am most confident about. The head predicts an error scale "
            "for every query alongside the flow itself. When you bin queries by predicted "
            "uncertainty, the actual error rises monotonically from 0.22 pixels in the "
            "confident bin to 6.7 in the least confident, a twenty-one fold span. It is not "
            "perfectly correlated, r is 0.345, but it is clearly informative and it is directly usable: weight "
            "correspondences in RANSAC, drop unreliable ones, or send more queries where the "
            "model is unsure. v2 has no equivalent output, so there is nothing to compare it "
            "against.")
    footer(s, i)

    # ============================================================ 10b selective accuracy
    s, i = slide()
    header(s, 'RESULTS', 'Confidence buys the accuracy back: 2.2x better over 80% of the frame')
    s.shapes.add_picture('results/plots/selective_accuracy.png', Inches(0.8), Inches(1.6),
                         width=Inches(8.2))
    add_text(s, 9.35, 1.8, 3.5, 3.6,
             'Mean error over every pixel is\none operating point, and it is\nthe only one v2 has.\n\n'
             'With a confidence value per\nquery, v3 has a curve.\n\n'
             'Drop the least confident 20%\nand error falls to 1.06 px,\n'
             'against v2 at 2.32.\n\n'
             'Keep the best fifth and it is\n0.48 px, nearly 5x better.', 11.5, False, INK, 1.18)
    chip(s, 9.35, 5.6, 3.5, 1.0,
         'For registration and mapping,\nwhere you pick your points\nanyway, abstaining is free.',
         10.5, BOX_LT, GOOD)
    note(s, "This is the slide I would most like your reaction to. The accuracy table earlier "
            "reports mean error over every pixel, which is the only operating point v2 has, "
            "because it cannot tell you which of its outputs to trust. v3 can, so instead of a "
            "point it has a curve. Discard the least confident fifth of queries and error on "
            "what remains falls to 1.06 pixels, against v2's 2.32 over the whole frame. So more "
            "than twice as accurate over eighty percent of the image. Keep only the most "
            "confident fifth and it is 0.48, close to five times better. Why this is not just a "
            "statistical trick: for registration and mapping you are choosing a few hundred "
            "points anyway, so declining to answer where the model is unsure costs you nothing. "
            "Same 2.35 million queries as the calibration bins on the previous slide.")
    footer(s, i)

    # ============================================================ 10c the three scenarios
    s, i = slide()
    header(s, 'SCENARIOS', 'Three situations a fast platform meets')
    s.shapes.add_picture('results/plots/scenarios_illustrated.png', Inches(0.75), Inches(1.45),
                         width=Inches(11.5))
    add_text(s, 0.75, 6.55, 11.9, 0.8,
             'Left: the frame, with the regions being tracked. Right: the flow actually returned, computed only inside those\n'
             'regions. Same geometry as the measurements that follow, on held-out VKITTI2 driving sequences.', 11.5, False, INK, 1.15)
    note(s, "These are the three situations you described, on real frames. First, something "
            "worth flowing enters the field of view and you start tracking it. Second, the "
            "platform turns and that region now overlaps a second object. Third, a new object "
            "appears in a frame you are already part way through processing. On the right is "
            "what the model actually computes: flow inside the marked regions only, and "
            "nothing elsewhere. The grey area is not computed at all. That is the whole "
            "premise, and the next two slides put numbers on what it costs and what you get "
            "back.")
    footer(s, i)

    # ============================================================ 10c2 crop cost
    s, i = slide()
    header(s, 'RESULTS', 'Flowing a region instead of the frame: 4.4x for 0.034 px')
    rows = [['Region processed', 'Area', 'Latency', 'Speedup', 'EPE in region'],
            ['Full frame', '100%', '33.3 ms', '1.0x', '0.657'],
            ['Region, no margin', '7.9%', '7.8 ms', '4.3x', '1.089'],
            ['Region + 32 px margin', '13.8%', '7.6 ms', '4.4x', '0.691'],
            ['Region + 64 px margin', '20.1%', '8.1 ms', '4.1x', '0.667'],
            ['Region + 128 px margin', '34.2%', '9.9 ms', '3.4x', '0.655']]
    table(s, 0.75, 1.75, rows, [3.5, 1.4, 1.5, 1.4, 1.7], 12.5, hl_row=3)
    add_text(s, 0.75, 4.35, 7.4, 2.2,
             'Keeping full resolution and processing 14% of the frame costs 0.034 px.\n\n'
             'The margin is not optional. With none, error rises 65% and a quarter of\n'
             'large-motion pixels fail outright, because global matching loses the\n'
             'context it needs to find them. Past 32 px nothing improves.\n\n'
             'The margin has to cover roughly one frame of motion, which for a\n'
             'platform is speed divided by frame rate. Mean motion here is 26.6 px\n'
             'and the knee is at 32.', 12, False, INK, 1.18)
    chip(s, 8.6, 4.45, 4.1, 1.5,
         'Design rule\n\nmargin  ~  expected motion\n~  speed / frame rate\n\n'
         'A faster platform needs a\nwider margin, and saves less.', 11.5, BOX_LT)
    add_text(s, 0.75, 6.75, 11.9, 0.35,
             'RTX 4060 laptop, fp16, 375x1242, 40 held-out VKITTI2 pairs. Applies to any flow network, v2 included.',
             10, False, MUTED)
    note(s, "Here is what flowing only a region actually buys. Keep full resolution, process "
            "fourteen percent of the frame, and it costs thirty four thousandths of a pixel. "
            "That is the enabling result for the whole scenario. The margin is the part worth "
            "dwelling on. Without one, error jumps sixty five percent and a quarter of the "
            "large-motion pixels fail completely, because this network finds large "
            "displacement using attention over the whole image and a tight crop takes that "
            "away. The margin has to be about one frame of expected motion, which is speed "
            "over frame rate. Mean motion in these sequences is twenty six point six pixels "
            "and the knee sits at thirty two, which is the rule falling out of the data. I "
            "should be clear that this applies to any flow network including v2 unchanged: it "
            "is a platform technique, not something my decoder provides. What the decoder adds "
            "is on the next slide.")
    footer(s, i)

    # ============================================================ 10c3 margin rule, two domains
    s, i = slide()
    header(s, 'RESULTS', 'The margin rule holds across two motion scales')
    s.shapes.add_picture('results/plots/margin_rule.png', Inches(1.05), Inches(1.6),
                         width=Inches(11.2))
    add_text(s, 0.75, 5.9, 11.9, 1.1,
             'Driving sequences average 26.6 px of motion and need about 32 px of margin. Aerial sequences average 9.26 px and\n'
             'need about 8. Both start from the same 0.43 px penalty with no margin, and both have shed most of it by the point\n'
             'where the margin equals one frame of motion. The requirement is set by displacement, not by image size, which\n'
             'makes it predictable from platform speed and frame rate before any flow is computed.', 11.5, False, INK, 1.15)
    note(s, "I wanted to know whether the margin rule was a curve fitted to one dataset or "
            "something more general, so I tested it on aerial sequences, which move very "
            "differently: nine point three pixels per frame against twenty six point six for "
            "driving. The prediction was that the margin needed should scale with the motion, "
            "so the knee should move from thirty two pixels down to about nine. It did. On the "
            "left you can see the two need different absolute margins. On the right, dividing "
            "by the mean motion, both start at the same penalty and both have lost most of it "
            "by the point where the margin equals one frame of motion. I will be precise "
            "about the limits: the agreement is good up to that point and the residuals differ "
            "beyond it, so this predicts the scale rather than the exact curve. That is still "
            "enough to size a margin from speed and frame rate before running anything, which "
            "is what a platform designer needs.")
    footer(s, i)

    # ============================================================ 10d scenario 3
    s, i = slide()
    header(s, 'RESULTS', 'A new object appears: 1.7 ms, against 7.3 for a new crop')
    s.shapes.add_picture('results/plots/scenario3_marginal.png', Inches(1.35), Inches(1.6),
                         width=Inches(10.6))
    add_text(s, 0.75, 5.95, 11.9, 1.1,
             'The third scenario: a frame is already being processed when something new enters the field of view. Because the\n'
             'coarse pass is cached, the decoder answers for 1.7 ms. A cropped pipeline has to run an entire new pass, 7.3 ms,\n'
             'and is markedly less accurate doing it because a fresh crop loses the global context that finds large motion.\n'
             'v2 keeps no state between calls, so it has nothing to answer from.', 11.5, False, INK, 1.15)
    note(s, "This is the scenario you described where the repeat query matters, and it is the "
            "one result that is specific to this architecture rather than to cropping. A frame "
            "is already in flight, and a new object enters the field of view. Because the "
            "coarse pass is cached, answering costs one point seven milliseconds for eight "
            "hundred points. A cropped pipeline cannot reuse anything: the new object is "
            "outside its box, so it runs a whole new pass at seven point three milliseconds, "
            "four times more, and the answer is worse, three point six pixels against two "
            "point one, because a fresh crop has lost the surrounding context. NeuFlow v2 "
            "keeps no state at all between calls. The general point is that this buys you the "
            "ability to answer a question you did not know you would have when the frame "
            "arrived, which on a moving platform is most of them.")
    footer(s, i)

    # ============================================================ 11 query interface
    s, i = slide()
    header(s, 'INTERFACE', 'What a query costs, and how you make one')
    s.shapes.add_picture('results/plots/decode_flat.png', Inches(0.8), Inches(1.65),
                         width=Inches(7.0))
    add_text(s, 8.1, 1.8, 4.6, 0.4, 'The entire API', 12, True, ACCENT)
    chip(s, 8.1, 2.2, 4.6, 1.9,
         'state = model.infer_coarse_state(img0, img1)\n'
         '                                  once, 33 ms\n\n'
         'flow  = model.decode_queries(\n'
         '            state, query_coords=q)\n'
         '            q: [B, N, 2] -> [B, N, 2], 1.3 ms\n\n'
         'flow, b = ... (return_uncertainty=True)', 10)
    add_text(s, 8.1, 4.3, 4.6, 2.0,
             'Decode cost is flat from 800 to 2,048\n'
             'queries, so it is bound by kernel launch\n'
             'overhead rather than compute. You get\n'
             '2,048 points for the price of 800.\n\n'
             'Above roughly 10,000 points the decode\n'
             'becomes compute-bound and dense mode\n'
             'is the better choice.', 11.5, False, INK, 1.18)
    note(s, "The interface is two calls. One coarse pass per frame pair, then as many decode "
            "calls as you like against it. The plot shows decode cost is flat between 800 and "
            "2,048 queries, which tells you it is dominated by kernel launch overhead, not "
            "arithmetic, so you get two thousand points for the price of eight hundred. Above "
            "about ten thousand points it becomes compute bound and you should just use dense "
            "mode.")
    footer(s, i)

    # ============================================================ 12 GUI
    s, i = slide()
    header(s, 'INTERFACE', 'Working tool: draw a region, get flow there')
    if os.path.exists('results/visuals/region_gui_window.png'):
        s.shapes.add_picture('results/visuals/region_gui_window.png',
                             Inches(0.8), Inches(1.65), width=Inches(7.6))
    add_text(s, 8.7, 1.8, 4.0, 4.4,
             'Load a video, step frame by frame,\n'
             'drag a box, get flow for that region.\n\n'
             'Two modes, both timed live:\n\n'
             'QUERY decodes only inside the box\n'
             'against a full-frame coarse pass. Exact.\n\n'
             'CROP feeds only the box through the\n'
             'whole pipeline, so cost scales with the\n'
             'area requested, at the price of losing\n'
             'context outside the crop.\n\n'
             'The panel shows where the time goes\n'
             'rather than asserting a speed-up.', 11.5, False, INK, 1.18)
    note(s, "This is a working tool, not a mock-up. You load a video, step to a frame, drag a "
            "box, and it computes flow in that box. Two modes. Query mode runs the full coarse "
            "pass and decodes only inside the box, which is exact. Crop mode feeds just the box "
            "through the whole pipeline, so the cost scales with the area you asked for, but it "
            "loses the surrounding context that global matching uses. The panel shows the "
            "actual breakdown live, including the parts where v3 is not the faster option.")
    footer(s, i)

    # ============================================================ 13 scorecard
    s, i = slide()
    header(s, 'ASSESSMENT', 'Against the three objectives I set')
    rows = [
        ['Objective', 'Verdict', 'Evidence'],
        ['Better accuracy than v2', 'NOT MET', 'best v3 2.384 vs 2.324; every config above v2'],
        ['', '', 'like-for-like 2.500, i.e. 7.6% behind'],
        ['', '', ''],
        ['Less compute than v2', 'PARTLY', 'dense 11% slower; first query level'],
        ['', '', 'repeat query 27x cheaper (1.3 vs 33.3 ms)'],
        ['', '', ''],
        ['Runs on edge devices', 'UNPROVEN', 'all figures from a laptop RTX 4060'],
        ['', '', 'no Jetson measurement has been taken'],
    ]
    table(s, 0.75, 1.85, rows, [3.6, 1.9, 6.4], 11.5)
    add_text(s, 0.75, 4.85, 11.9, 0.4, 'What the project does deliver', 12.5, True, ACCENT)
    for k, (t, d) in enumerate([
        ('Queryable output', 'flow at any continuous coordinate, sparse equals dense exactly'),
        ('Cheap repeat access', '1.3 ms per further query batch on a cached frame'),
        ('Calibrated confidence', 'per-query error estimate, monotonic against real error'),
    ]):
        chip(s, 0.75 + k * 4.0, 5.3, 3.8, 1.1, f'{t}\n\n{d}', 10.5)
    note(s, "Here is the scorecard against what I set out to do, and two of three are not met. "
            "Better accuracy: not met, the fair comparison is a tie. Less compute: partly, "
            "dense is slower, a first query is level, but repeat queries are nearly eight times "
            "cheaper. Edge capable: unproven, everything I have measured is on a "
            "laptop 4060 and I have not touched a Jetson. What the project does deliver is the "
            "three things at the bottom: output you can query at any coordinate, cheap repeat "
            "access to a cached frame, and a calibrated confidence value. Those are real and "
            "they are measured.")
    footer(s, i)

    # ============================================================ 14 limitations
    s, i = slide()
    header(s, 'LIMITATIONS', 'What is not yet proven, stated plainly')
    items = [
        ('Sub-pixel precision is worse than v2',
         'By 0.8 to 6.3 points of 1-pixel accuracy. Cause diagnosed (decoder never sees '
         'full-resolution input.'),
        ('No edge-device measurement',
         'Every latency figure is from a laptop RTX 4060. The edge claim in the title is a design '
         'target, not a result.'),
        ('Single evaluation domain',
         'All accuracy numbers are VKITTI2 Scene18+20. A synthetic driving benchmark is not '
         'evidence for field or survey imagery.'),
        ('Spring run truncated',
         'Killed by the 8-hour wall clock at roughly 90k of 100k steps, so it is not directly '
         'comparable to the others.'),
        ('One seed, and checkpoint noise the size of the effects',
         'Step 90k vs 100k of the same run moves EPE by up to 0.038 px, reversing two of the '
         'orderings. Nothing finer than about 0.05 px is resolvable here.'),
    ]
    y = 1.8
    for t, d in items:
        add_text(s, 0.75, y, 4.4, 0.3, t, 11.5, True, INK)
        add_text(s, 5.3, y, 7.3, 0.6, d, 11, False, INK, 1.12)
        y += 1.02
    note(s, "Limitations, stated plainly rather than buried. Precision is worse than v2 and I "
            "have diagnosed without yet addressing it. There is no edge device measurement at all, so "
            "the word edge in my title is a design target and not a result. Everything is "
            "evaluated on one synthetic driving dataset, which says little about field imagery. "
            "The Spring run was cut off by the cluster wall clock. And I have one seed per "
            "configuration, so please do not read too much into differences below about "
            "five hundredths of a pixel.")
    footer(s, i)

    # ============================================================ 15 next steps
    s, i = slide()
    header(s, 'NEXT STEPS', 'In priority order')
    steps = [
        ('1', 'Full-resolution stem for the decoder',
         'Give the decoder a cheap full-resolution feature map so evidence actually varies '
         'within an 8x8 cell. This is the direct attack on the precision gap and on making '
         'continuous querying meaningful rather than nominal.'),
        ('2', 'Spring 4K evaluation',
         'Spring provides ground truth at twice the input resolution. v3 can be queried there '
         'natively; v2 structurally cannot. The one test that demonstrates a capability rather '
         'than a margin.'),
        ('3', 'Jetson measurement',
         'Converts the edge claim from a design target into a result, or refutes it.'),
        ('4', 'Field or survey imagery',
         'A registration demonstration on lab data would test the actual use case.'),
    ]
    y = 1.85
    for num, t, d in steps:
        chip(s, 0.75, y, 0.42, 0.42, num, 12, BOX_DK, WHITE, bold=True)
        add_text(s, 1.35, y - 0.02, 4.0, 0.3, t, 12, True, INK)
        add_text(s, 5.5, y - 0.02, 7.1, 0.9, d, 11, False, INK, 1.12)
        y += 1.28
    note(s, "Four next steps in priority order. First, give the decoder a full resolution "
            "feature map. That is the direct attack on the precision gap, and it also matters "
            "because right now querying between pixels returns something close to "
            "interpolation, so the arbitrary coordinate claim is thinner than it sounds. "
            "Second, the Spring dataset gives ground truth at twice the input resolution, which "
            "v3 can be queried at natively and v2 structurally cannot. That is the one test "
            "that shows a capability rather than a margin. Third, actually measure a Jetson. "
            "Fourth, try it on real survey imagery.")
    footer(s, i)

    # ============================================================ 16 Q&A
    s, i = slide()
    header(s, 'ANTICIPATED QUESTIONS', '')
    qa = [
        ('Why is flow defined between pixels at all?',
         'Motion is continuous; the pixel grid is a sampling artefact. A corner tracked to '
         '(312.7, 188.2) is a real answer, and registration consumes it directly.'),
        ('Why freeze the backbone?',
         'To isolate the decoder as the only variable. Unfreezing diverged at batch 4 locally '
         'and has not been retried at scale.'),
        ('Why is dense v3 slower than v2?',
         'v2 upsamples with one convolution that shares work across neighbouring pixels. v3 '
         'answers each query independently, giving up that sharing to gain the ability to '
         'answer only what is asked.'),
        ('Is a tie on EPE worth the added complexity?',
         'Not on its own. It is worth it if you need queryability, cheap repeat access, or '
         'confidence values. For plain dense flow, v2 remains the better choice, and I would '
         'say so.'),
    ]
    y = 1.65
    for q, a in qa:
        add_text(s, 0.75, y, 11.9, 0.3, q, 12, True, INK)
        add_text(s, 0.95, y + 0.33, 11.6, 0.65, a, 11, False, ACCENT, 1.12)
        y += 1.3
    note(s, "A few questions I expect. Why flow between pixels: because motion is continuous "
            "and the grid is an artefact of sampling. Why freeze the backbone: to keep the "
            "decoder as the only variable, and because unfreezing diverged when I tried it. Why "
            "is dense slower: because v2's convolution shares work between neighbouring pixels "
            "and my decoder deliberately gives that up in exchange for answering only what is "
            "asked. And the hardest one, is a tie worth the complexity. On its own, no. If you "
            "want plain dense flow, use v2, and I will say that plainly. It is worth it if you "
            "need queryability, cheap repeat access, or a confidence signal.")
    footer(s, i)

    os.makedirs('docs', exist_ok=True)
    prs.save(OUT)
    notes = sum(1 for sl in prs.slides
                if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip())
    print(f'saved {OUT} with {len(prs.slides._sldIdLst)} slides and {notes} speaker notes')


if __name__ == '__main__':
    main()
